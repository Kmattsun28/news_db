import feedparser
import os
import requests
import pdfplumber
from pptx import Presentation
from datetime import datetime, timedelta
from app.script.db import SessionLocal
from app.script.models import NewsArticle
from app.script.debug import debug_printer as d
from app.script.utils_scraper import extract_article_text, detect_currency_tags
from app.script.summarizer import summarize_text
# Finnhub APIの追加
from app.script.finnhub_news import fetch_finnhub_forex_news

# RSS_FEEDSの定義は get_optimized_rss_feeds() 関数内に移動しました



def is_recent_article(published_parsed, hours_back=24):
    """
    記事が指定した時間以内に公開されたかどうかをチェック
    
    Args:
        published_parsed: feedparserのpublished_parsedフィールド
        hours_back: 何時間前までの記事を有効とするか（デフォルト: 24時間）
    
    Returns:
        bool: 指定時間以内の記事の場合True、それ以外はFalse
    """
    if not published_parsed:
        return False
    
    try:
        article_time = datetime(*published_parsed[:6])
        cutoff_time = datetime.now() - timedelta(hours=hours_back)
        is_recent = article_time >= cutoff_time
        
        # if not is_recent:
        #     d.print(f"⏰ Article too old: {article_time.strftime('%Y-%m-%d %H:%M:%S')} (cutoff: {cutoff_time.strftime('%Y-%m-%d %H:%M:%S')})", level="debug")
        
        return is_recent
    except Exception as e:
        d.print(f"Error checking article date: {e}", level="warning")
        return False


def extract_pdf_text(url: str) -> str:
    try:
        response = requests.get(url)
        with open("/tmp/tmp.pdf", "wb") as f:
            f.write(response.content)
        text = ""
        with pdfplumber.open("/tmp/tmp.pdf") as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        d.print(f"PDF抽出エラー: {e}")
        return ""

def extract_pptx_text(url: str) -> str:
    try:
        response = requests.get(url)
        with open("/tmp/tmp.pptx", "wb") as f:
            f.write(response.content)
        prs = Presentation("/tmp/tmp.pptx")
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        d.print(f"PPTX抽出エラー: {e}")
        return ""
    
def extract_xlsx_text(url: str) -> str:
    import subprocess
    tmp_xlsx = "/tmp/tmp.xlsx"
    tmp_pdf = "/tmp/tmp_xlsx.pdf"
    try:
        response = requests.get(url)
        with open(tmp_xlsx, "wb") as f:
            f.write(response.content)
        # LibreOfficeでxlsx→pdf変換
        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf", "--outdir", "/tmp", tmp_xlsx
        ], check=True)
        # PDFからテキスト抽出
        text = ""
        if os.path.exists(tmp_pdf):
            with pdfplumber.open(tmp_pdf) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        d.print(f"XLSX抽出エラー: {e}")
        return ""
    finally:
        for path in [tmp_xlsx, tmp_pdf]:
            if os.path.exists(path):
                os.remove(path)

def fetch_and_store_rss():
    session = SessionLocal()
    d.print_ts(f"<<< scheduled task: fetch_and_store_rss >>>", level='debug')
    
    total_processed = 0
    total_added = 0
    batch_size = 10  # バッチサイズを設定
    current_batch = 0

    try:
        # 1. 時間フィルタリング対応フィードを先に処理（効率的）
        time_filtered_feeds, standard_feeds = get_optimized_rss_feeds()
        
        # 時間指定可能なフィードを処理（Google Newsなど）
        d.print("Processing time-filtered RSS feeds...", level="info")
        for category, urls in time_filtered_feeds.items():
            d.print(f"Fetching time-filtered RSS feeds for category: {category} contains {len(urls)}", level="debug", output_path="./data/fetch_and_store_rss.log")
            for url in urls:
                try:
                    # キャッシュ情報を取得
                    cache_info = get_feed_cache_info(url)
                    
                    # 効率的RSS取得（etag/modified使用）
                    feed, new_etag, new_modified = fetch_rss_with_caching(
                        url, 
                        cache_info['etag'], 
                        cache_info['modified']
                    )
                    
                    if feed is None:  # 304 Not Modified or error
                        continue
                    
                    # キャッシュ情報を更新
                    save_feed_cache_info(url, new_etag, new_modified)
                    
                    d.print(f"Processing time-filtered feed: {url}, entries: {len(feed.entries)}", level="debug")
                    
                    # 時間フィルタ済みのため、全記事を処理
                    for entry in feed.entries:
                        try:
                            total_processed += 1
                            
                            published = datetime(*entry.published_parsed[:6]) if entry.get("published_parsed") else datetime.now()
                            exists = session.query(NewsArticle).filter_by(title=entry.title, published=published).first()
                            if exists:
                                d.print(f"⏩ skip article: {entry.title[:50]}... (already exists)", output_path="./data/fetch_and_store_rss.log")
                                continue

                            # 記事処理（既存のロジック）
                            ext = os.path.splitext(entry.link)[1].lower()
                            full_text = ""
                            
                            try:
                                if ext == ".pdf":
                                    d.print(f"Processing PDF: {entry.link}", level="warning")
                                    full_text = extract_pdf_text(entry.link)
                                elif ext in [".ppt", ".pptx"]:
                                    d.print(f"Processing PPTX: {entry.link}", level="warning")
                                    full_text = extract_pptx_text(entry.link)
                                elif ext in [".xls", ".xlsx"]:
                                    d.print(f"Processing XLSX: {entry.link}", level="warning")
                                    full_text = extract_xlsx_text(entry.link)
                                else:
                                    full_text = extract_article_text(entry.link)
                            except Exception as scrape_error:
                                d.print(f"Scraping failed for {entry.link}: {scrape_error}", level="warning")
                                full_text = ""

                            if not full_text:
                                full_text = entry.get("summary", entry.get("title", ""))

                            try:
                                summary = summarize_text(full_text)
                                currency_tags = detect_currency_tags(full_text)
                            except Exception as ai_error:
                                d.print(f"AI processing failed for {entry.title[:50]}...: {ai_error}", level="warning")
                                summary = entry.get("summary", entry.get("title", ""))
                                currency_tags = []

                            article = NewsArticle(
                                category=category,
                                title=entry.title,
                                summary=summary,
                                url=entry.link,
                                published=published,
                                currency_tags=currency_tags
                            )
                            session.add(article)
                            total_added += 1
                            current_batch += 1
                            d.print(f"✅ Added time-filtered RSS article: {entry.title[:50]}...", level="debug")
                            
                            # バッチサイズに達したら中間コミット
                            if current_batch >= batch_size:
                                try:
                                    session.commit()
                                    d.print(f"📦 Batch commit: {current_batch} articles saved", level="info")
                                    current_batch = 0
                                except Exception as commit_error:
                                    d.print(f"❌ Batch commit failed: {commit_error}", level="error")
                                    session.rollback()
                                    current_batch = 0
                            
                        except Exception as entry_error:
                            d.print(f"Error processing entry {entry.get('title', 'Unknown')[:50]}...: {entry_error}", level="error")
                            continue
                            
                except Exception as feed_error:
                    d.print(f"Error processing time-filtered feed {url}: {feed_error}", level="error")
                    continue
        
        # 2. 標準フィード（時間指定不可）を処理
        d.print("Processing standard RSS feeds with post-filtering...", level="info")
        for category, urls in standard_feeds.items():
            d.print(f"Fetching standard RSS feeds for category: {category} contains {len(urls)}", level="debug", output_path="./data/fetch_and_store_rss.log")
            for url in urls:
                try:
                    # 標準的なRSS取得
                    feed = feedparser.parse(url)
                    d.print(f"Processing standard feed: {url}, entries: {len(feed.entries)}", level="debug")
                    
                    for entry in feed.entries:
                        try:
                            total_processed += 1
                            
                            # 記事の公開日をチェック（過去24時間以内のもののみ処理）
                            if not is_recent_article(entry.get("published_parsed"), hours_back=24):
                                # d.print(f"⏰ Skipping old article: {entry.title[:50]}... (published too long ago)", level="debug", output_path="./data/fetch_and_store_rss.log")
                                continue
                            
                            published = datetime(*entry.published_parsed[:6]) if entry.get("published_parsed") else datetime.now()
                            exists = session.query(NewsArticle).filter_by(title=entry.title, published=published).first()
                            if exists:
                                d.print(f"⏩ skip article: {entry.title[:50]}... (already exists)", output_path="./data/fetch_and_store_rss.log")
                                continue

                            # RSSエントリのリンクの種類で処理を分岐 http, pdf, pptx, xlsx
                            ext = os.path.splitext(entry.link)[1].lower()
                            full_text = ""
                            
                            try:
                                if ext == ".pdf":
                                    d.print(f"Processing PDF: {entry.link}", level="warning")
                                    full_text = extract_pdf_text(entry.link)
                                elif ext in [".ppt", ".pptx"]:
                                    d.print(f"Processing PPTX: {entry.link}", level="warning")
                                    full_text = extract_pptx_text(entry.link)
                                elif ext in [".xls", ".xlsx"]:
                                    d.print(f"Processing XLSX: {entry.link}", level="warning")
                                    full_text = extract_xlsx_text(entry.link)
                                else:
                                    full_text = extract_article_text(entry.link)
                            except Exception as scrape_error:
                                d.print(f"Scraping failed for {entry.link}: {scrape_error}", level="warning")
                                full_text = ""

                            if not full_text:
                                full_text = entry.get("summary", entry.get("title", ""))

                            try:
                                summary = summarize_text(full_text)
                                currency_tags = detect_currency_tags(full_text)
                            except Exception as ai_error:
                                d.print(f"AI processing failed for {entry.title[:50]}...: {ai_error}", level="warning")
                                summary = entry.get("summary", entry.get("title", ""))
                                currency_tags = []

                            article = NewsArticle(
                                category=category,
                                title=entry.title,
                                summary=summary,
                                url=entry.link,
                                published=published,
                                currency_tags=currency_tags
                            )
                            session.add(article)
                            total_added += 1
                            current_batch += 1
                            d.print(f"✅ Added standard RSS article: {entry.title[:50]}...", level="debug")
                            
                            # バッチサイズに達したら中間コミット
                            if current_batch >= batch_size:
                                try:
                                    session.commit()
                                    d.print(f"📦 Batch commit: {current_batch} articles saved", level="info")
                                    current_batch = 0
                                except Exception as commit_error:
                                    d.print(f"❌ Batch commit failed: {commit_error}", level="error")
                                    session.rollback()
                                    current_batch = 0
                            
                        except Exception as entry_error:
                            d.print(f"Error processing standard entry {entry.get('title', 'Unknown')[:50]}...: {entry_error}", level="error")
                            continue
                            
                except Exception as feed_error:
                    d.print(f"Error processing standard feed {url}: {feed_error}", level="error")
                    continue

        # 残りの記事をコミット
        if current_batch > 0:
            try:
                session.commit()
                d.print(f"📦 Final batch commit: {current_batch} articles saved", level="info")
            except Exception as commit_error:
                d.print(f"❌ Final batch commit failed: {commit_error}", level="error")
                session.rollback()

        d.print(f"🆗 All RSS feeds processed. Total processed: {total_processed}, Added: {total_added}", level="info")
        
    except Exception as e:
        session.rollback()
        d.print(f"Error in fetch_and_store_rss: {e}", level="error")
    finally:
        session.close()


def fetch_and_store_all_news():
    """
    すべてのニュースソース（RSS + Finnhub API）からニュースを取得してDBに保存
    
    RSS収集期間: 過去24時間以内の記事のみ
    Finnhub収集期間: 過去30分間
    """
    d.print_ts("<<< scheduled task: fetch_and_store_all_news >>>", level='info')
    
    total_new_articles = 0
    
    # 1. 既存のRSSフィードからニュース取得（過去24時間以内）
    d.print("=== RSS News Collection (Last 24 hours) ===", level="info")
    try:
        fetch_and_store_rss()
        d.print("✅ RSS news collection completed", level="info")
    except Exception as e:
        d.print(f"❌ RSS news collection failed: {e}", level="error")
    
    # 2. Finnhub APIからニュース取得（過去30分間）
    d.print("=== Finnhub News Collection (Last 30 minutes) ===", level="info")
    try:
        finnhub_count = fetch_finnhub_forex_news(limit=50, minutes_back=30)
        total_new_articles += finnhub_count
        d.print(f"✅ Finnhub news collection completed. New articles: {finnhub_count}", level="info")
    except Exception as e:
        d.print(f"❌ Finnhub news collection failed: {e}", level="error")
    
    d.print(f"🏁 All news collection completed. Total new articles from Finnhub: {total_new_articles}", level="info")
    return total_new_articles


def generate_time_filtered_rss_urls():
    """
    時間範囲を指定したRSS URLを動的に生成
    Google Newsなど一部のRSSフィードは時間範囲パラメータに対応
    """
    from datetime import datetime, timedelta
    
    # 24時間前の日時を取得（開始日の前日）
    start_time = datetime.now() - timedelta(hours=24)
    after_param = (start_time - timedelta(days=1)).strftime("%Y-%m-%d")  # 開始日の前日
    before_param = (datetime.now() + timedelta(days=1)).strftime("%Y-%m-%d")  # 終了日の翌日
    
    # 時間指定可能なRSSフィード（Google News）
    time_filtered_feeds = {
        "google_news": [
            f"https://news.google.com/rss/search?q=USD+JPY+exchange+rate+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=EUR+USD+exchange+rate+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=Federal+Reserve+rate+decision+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=US+inflation+CPI+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=nonfarm+payrolls+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=ECB+Lagarde+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=EUR+inflation+after:{after_param}+before:{before_param}&hl=en-US&gl=US&ceid=US:en",
            f"https://news.google.com/rss/search?q=日銀+為替介入+after:{after_param}+before:{before_param}&hl=ja&gl=JP&ceid=JP:ja",
            f"https://news.google.com/rss/search?q=円安+after:{after_param}+before:{before_param}&hl=ja&gl=JP&ceid=JP:ja",
            f"https://news.google.com/rss/search?q=日米金利差+after:{after_param}+before:{before_param}&hl=ja&gl=JP&ceid=JP:ja",
        ]
    }
    
    return time_filtered_feeds


def get_optimized_rss_feeds():
    """
    時間フィルタリング対応・非対応を分けてRSSフィードを取得
    """
    # 時間指定可能なフィード
    time_filtered = generate_time_filtered_rss_urls()
    
    # 時間指定不可のフィード（従来通り後でフィルタ）
    standard_feeds = {
        # -------------------------------------------------------------------
        # 中央銀行・政府系機関（為替・政策発表）
        # -------------------------------------------------------------------
        "official_sources": [
            "https://www.boj.or.jp/rss/whatsnew.rdf",                          # 日本銀行
            "https://www.mof.go.jp/public_relations/rss.xml",                 # 財務省（為替介入含む）
            "https://www.federalreserve.gov/feeds/pressreleases.xml",         # FRB
            "https://home.treasury.gov/news/press-releases/feed",             # 米財務省
            "https://www.ecb.europa.eu/press/rss/pr.xml",                     # 欧州中央銀行
            # "https://www.ecb.europa.eu/stats/eurofxref/eurofxref-usd.rss",    # EUR/USD為替レート（ECB公式）
            # "https://www.bundesbank.de/rss/bbk_en_news.xml",                  # ドイツ連邦銀行（英語ニュース）
        ],
        
        # -------------------------------------------------------------------
        # 専門FXニュース/分析（USD, JPY, EUR中心）
        # -------------------------------------------------------------------
        "forex_feeds": [
            # "https://www.forexlive.com/feed/",                                 # ForexLive（即時ニュース）
            # "https://www.fxstreet.com/rss/news",                               # FXStreet - ニュース
            # "https://www.fxstreet.com/rss/analysis",                           # FXStreet - 分析
            # "https://www.fxstreet.com/rss/forex-news",                         # FXStreet - 為替限定
            # "https://www.fxstreet.com/rss/forex-technical-analysis",           # FXStreet - テクニカル分析
            # "https://www.investing.com/rss/news_forex.xml",                   # Investing.com - FXニュース
            # "https://www.investing.com/rss/technical.xml",                    # Investing.com - テクニカル
            # "https://www.actionforex.com/feed",                                # ActionForex
        ],
        
        # -------------------------------------------------------------------
        # 国際経済・金融報道（USD・EUR政策に影響）
        # -------------------------------------------------------------------
        "international_media_en": [
            # "https://feeds.bloomberg.com/markets/news.rss",                   # Bloomberg - Markets
            # "https://feeds.a.dj.com/rss/RSSMarketsMain.xml",                 # Wall Street Journal - Markets
            # "https://rss.nytimes.com/services/xml/rss/nyt/Business.xml",     # NYTimes - Business
            "https://www.economist.com/finance-and-economics/rss.xml",       # The Economist - Finance
        ],
        
        # -------------------------------------------------------------------
        # 国内メディア（為替/金融/経済）
        # -------------------------------------------------------------------
        "domestic_media_jp": [
            "https://www.nikkei.com/rss/newselement/Nni01.xml",              # 日経 - 主要ニュース
            # "https://www.nikkei.com/rss/market/market.xml",                  # 日経 - マーケット
            # "https://toyokeizai.net/list/feed/rss",                          # 東洋経済 - 総合RSS
            # "https://www3.nhk.or.jp/rss/news/cat0.xml",                      # NHK - 総合経済
        ]
    }
    
    return time_filtered, standard_feeds


def fetch_rss_with_caching(url, etag=None, modified=None):
    """
    RSS取得時にetagとmodifiedヘッダーを使用して効率的に取得
    前回取得時から変更がない場合は304 Not Modifiedが返される
    
    Args:
        url: RSS URL
        etag: 前回取得時のetag
        modified: 前回取得時のmodified時刻
    
    Returns:
        tuple: (feed, new_etag, new_modified)
    """
    try:
        # etagとmodifiedがある場合、それらを指定してリクエスト
        feed = feedparser.parse(url, etag=etag, modified=modified)
        
        # 304 Not Modified の場合、新しい記事なし
        if feed.status == 304:
            d.print(f"📄 No new articles in feed: {url}", level="debug")
            return None, etag, modified
        
        # 新しいetagとmodifiedを保存
        new_etag = feed.get('etag')
        new_modified = feed.get('modified_parsed')
        
        return feed, new_etag, new_modified
        
    except Exception as e:
        d.print(f"Error fetching RSS with caching {url}: {e}", level="error")
        return None, None, None


def get_feed_cache_info(url):
    """
    RSSフィードのキャッシュ情報を取得（実装例）
    実際の実装では、データベースやファイルに保存
    """
    # 簡易実装：メモリベース（実用では永続化が必要）
    cache_store = getattr(get_feed_cache_info, 'cache', {})
    return cache_store.get(url, {'etag': None, 'modified': None})


def save_feed_cache_info(url, etag, modified):
    """
    RSSフィードのキャッシュ情報を保存
    """
    if not hasattr(save_feed_cache_info, 'cache'):
        save_feed_cache_info.cache = {}
    save_feed_cache_info.cache[url] = {'etag': etag, 'modified': modified}